import glob, io, json
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Emu
import img2pdf

OUT = Path("out"); slides = sorted(glob.glob("out/slides/slide-*.png"))
W, H = 6858000, 12192000   # 7.5in x 13.333in (9:16 portrait)

# --- PPTX: one full-bleed picture per slide, nothing else -------------------
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(W), Emu(H)
blank = prs.slide_layouts[6]
for f in slides:
    s = prs.slides.add_slide(blank)
    # strip any placeholder shapes that a layout might carry
    for shp in list(s.shapes):
        shp._element.getparent().remove(shp._element)
    s.shapes.add_picture(f, 0, 0, width=Emu(W), height=Emu(H))
pptx_path = OUT / "Two_Coasts_LuxuryProperty_Phone_Safe.pptx"
prs.save(pptx_path)

# --- PDF for WhatsApp: high-quality JPEG pages, 9:16 ------------------------
jpgs = []
for f in slides:
    im = Image.open(f).convert("RGB")
    b = io.BytesIO(); im.save(b, "JPEG", quality=90, optimize=True, subsampling=0); jpgs.append(b.getvalue())
layout = img2pdf.get_layout_fun((img2pdf.in_to_pt(7.5), img2pdf.in_to_pt(13.3333)))
pdf_path = OUT / "Two_Coasts_LuxuryProperty_Phone_Safe.pdf"
pdf_path.write_bytes(img2pdf.convert(jpgs, layout_fun=layout))

# --- verify the PPTX: every slide exactly one picture, full-bleed -----------
p2 = Presentation(pptx_path)
ok = True
for i, s in enumerate(p2.slides, 1):
    shapes = list(s.shapes)
    good = len(shapes) == 1 and shapes[0].shape_type == 13 and shapes[0].left == 0 and shapes[0].top == 0 \
        and shapes[0].width == W and shapes[0].height == H
    ok &= good
    if not good: print("slide", i, [(sh.shape_type, sh.name) for sh in shapes])
print("pptx slides:", len(p2.slides), "size", p2.slide_width, p2.slide_height, "all single full-bleed image:", ok)
print("pptx bytes:", pptx_path.stat().st_size, "pdf bytes:", pdf_path.stat().st_size)
